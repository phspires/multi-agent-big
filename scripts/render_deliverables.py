#!/usr/bin/env python3
"""Render BIG Impact PoV Review deliverables.

Fills the .pptx and .docx templates with content from state/drafts/ and
state/evidence_ledger.json. Outputs go to outputs/.

Field mapping (PPTX) — placeholder shapes identified by shape_id (verified by
inspecting the unzipped slide XML before writing this script):

Slide 1:
  - Group 21 (Problem Statement large box)        -> problem.md Draft text
  - Group 6  (Hipótese de valor)                  -> value.md Main Hypothesis
  - Group 16 ((principal) métrica)                -> value.md Principal Metric (line)
  - TextBox 35, 36 ((suporte) x2)                  -> value.md Support Metric 1 / 2 (line)
  - Group 4  (Tabela de Evidências large box)     -> evidence_table.json rows
Slide 2:
  - Group 4  (Pilot Sketch large box)             -> pilot.md Pilot Design summary
  - Group 20 (Budget Range value box)             -> pilot.md TOTAL row
  - Group 25 (budget notes line below Budget)     -> pilot.md TOTAL notes
  - TextBox 39, 40 ((preencher) x2 in Pressupostos)-> pilot.md top 2 critical assumptions

Field mapping (DOCX): contact-info placeholders -> static project info; then
append "Notas Adicionais" content with sources scanned, per-claim notes for
EV-IDs cited in the deck, and a methodology paragraph.
"""

from __future__ import annotations

import copy
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from docx import Document
from docx.shared import Pt as DocxPt

ROOT = Path("/Users/pedro/git_repos/big_impact")
DRAFTS = ROOT / "state" / "drafts"
LEDGER_FILE = ROOT / "state" / "evidence_ledger.json"
PPTX_TEMPLATE = ROOT / "BIG Impact - Proof of Value Review Template.pptx"
DOCX_TEMPLATE = ROOT / "Big Impact - PoV Review - Evidence Notes Template.docx"
OUT_DIR = ROOT / "outputs"
PPTX_OUT = OUT_DIR / "PoV_Review_PerSense.pptx"
DOCX_OUT = OUT_DIR / "Evidence_Notes_PerSense.docx"

PROJECT_NAME = "PerSense"
CONTACT = "Pedro Pires — pedro.pires@gmail.com"

# Soft length budgets (raise warning, do NOT auto-truncate).
PROBLEM_MAX = 500          # chars (case rule)
HYPOTHESIS_MAX = 100       # chars (value.md self-imposed)
METRIC_LINE_MAX = 220
SUPPORT_LINE_MAX = 180
PILOT_SKETCH_MAX = 1100
ASSUMPTION_MAX = 220
BUDGET_LINE_MAX = 60
BUDGET_NOTES_MAX = 200
EVIDENCE_ROW_MAX = 220


# ---------------------------------------------------------------------------
# Loaders
# ---------------------------------------------------------------------------


def load_drafts() -> dict:
    problem_md = (DRAFTS / "problem.md").read_text(encoding="utf-8")
    value_md = (DRAFTS / "value.md").read_text(encoding="utf-8")
    pilot_md = (DRAFTS / "pilot.md").read_text(encoding="utf-8")
    evidence_table = json.loads((DRAFTS / "evidence_table.json").read_text(encoding="utf-8"))
    ledger = json.loads(LEDGER_FILE.read_text(encoding="utf-8"))
    return dict(
        problem_md=problem_md,
        value_md=value_md,
        pilot_md=pilot_md,
        evidence_table=evidence_table,
        ledger=ledger,
    )


# ---------------------------------------------------------------------------
# Draft parsers
# ---------------------------------------------------------------------------


def extract_problem(problem_md: str) -> str:
    # The Draft block is between **Draft** ... **Evidence References**
    m = re.search(r"\*\*Draft\*\*[^\n]*:\s*\n(.+?)\n\n\*\*Evidence References", problem_md, re.S)
    if not m:
        raise ValueError("Could not find Draft block in problem.md")
    return m.group(1).strip()


def extract_value(value_md: str) -> dict:
    # Main Hypothesis: line right after the header
    m = re.search(r"\*\*Main Hypothesis\*\*[^\n]*\n(.+?)\n\n", value_md, re.S)
    if not m:
        raise ValueError("Could not find Main Hypothesis")
    main_hypothesis = m.group(1).strip()

    def block(name: str) -> str:
        # Capture from "**name**:" up to next "**" header or end
        m = re.search(rf"\*\*{re.escape(name)}\*\*:\s*\n(.+?)(?=\n\*\*|\Z)", value_md, re.S)
        if not m:
            raise ValueError(f"Could not find block: {name}")
        return m.group(1).strip()

    def kv(block_text: str, key: str) -> str:
        m = re.search(rf"-\s*{re.escape(key)}:\s*(.+?)(?=\n-\s|\Z)", block_text, re.S)
        if not m:
            raise ValueError(f"Could not find key {key!r} in block")
        return " ".join(m.group(1).split()).strip()

    pm = block("Principal Metric")
    sm1 = block("Support Metric 1")
    sm2 = block("Support Metric 2")

    return {
        "main_hypothesis": main_hypothesis,
        "principal_metric_name": kv(pm, "Name"),
        "principal_metric_baseline": kv(pm, "Baseline"),
        "principal_metric_target": kv(pm, "Target"),
        "support1_name": kv(sm1, "Name"),
        "support1_hypothesis": kv(sm1, "Hypothesis"),
        "support2_name": kv(sm2, "Name"),
        "support2_hypothesis": kv(sm2, "Hypothesis"),
    }


def extract_pilot(pilot_md: str) -> dict:
    # Operational Hypothesis line
    m = re.search(r"\*\*Operational Hypothesis\*\*:\s*(.+?)(?=\n\n)", pilot_md, re.S)
    op_hypothesis = " ".join(m.group(1).split()).strip() if m else ""

    # Timeline rows for a compact text representation
    timeline_rows = re.findall(r"^\|\s*(–?[\d a]+?)\s*\|\s*(.+?)\s*\|\s*(.+?)\s*\|\s*$", pilot_md, re.M)
    # Filter out header/separator
    timeline_rows = [r for r in timeline_rows if not r[0].startswith("---") and r[0] != "Sem"]

    # TOTAL budget row
    m_total = re.search(r"\|\s*\*\*TOTAL\*\*\s*\|\s*\*\*(.+?)\*\*\s*\|\s*(.+?)\s*\|", pilot_md)
    if not m_total:
        raise ValueError("Could not find TOTAL budget row")
    budget_total = m_total.group(1).strip()
    budget_notes = m_total.group(2).strip()

    # Critical Assumptions: numbered list under the section header.
    sect = re.search(
        r"## Critical Assumptions[^\n]*\n+(.+?)(?=\n## )", pilot_md, re.S
    )
    assumption_blocks = []
    if sect:
        # Each item: "1. **Title**: text"
        for m_a in re.finditer(r"^\d+\.\s+(.+?)(?=\n\d+\.\s|\Z)", sect.group(1), re.S | re.M):
            assumption_blocks.append(" ".join(m_a.group(1).split()).strip())

    return {
        "operational_hypothesis": op_hypothesis,
        "timeline_rows": timeline_rows,
        "budget_total": budget_total,
        "budget_notes": budget_notes,
        "assumptions": assumption_blocks,
    }


# ---------------------------------------------------------------------------
# EV-ID extraction
# ---------------------------------------------------------------------------

EV_ID_RE = re.compile(r"\bEV-\d{3}\b")


def collect_evids_in_deck_text(parts: dict, evidence_rows: list[dict]) -> list[str]:
    """Collect EV-IDs that will end up rendered onto the slides.

    Only count EV-IDs that appear in the actual rendered text (not in markdown
    sections we don't render).
    """
    ids: set[str] = set()
    # Problem text
    ids.update(EV_ID_RE.findall(parts["problem_text"]))
    # Value: hypothesis + metrics lines
    v = parts["value"]
    for s in (
        v["main_hypothesis"],
        v["principal_metric_name"],
        v["principal_metric_baseline"],
        v["principal_metric_target"],
        v["support1_name"],
        v["support1_hypothesis"],
        v["support2_name"],
        v["support2_hypothesis"],
    ):
        ids.update(EV_ID_RE.findall(s))
    # Pilot
    p = parts["pilot"]
    ids.update(EV_ID_RE.findall(p["operational_hypothesis"]))
    ids.update(EV_ID_RE.findall(p["budget_total"]))
    ids.update(EV_ID_RE.findall(p["budget_notes"]))
    for a in p["assumptions"]:
        ids.update(EV_ID_RE.findall(a))
    # Evidence table rows
    for row in evidence_rows:
        ids.add(row["evidence_id"])
    return sorted(ids, key=lambda x: int(x.split("-")[1]))


# ---------------------------------------------------------------------------
# PPTX rendering helpers
# ---------------------------------------------------------------------------


def find_shape_by_id(shapes, shape_id):
    """Recurse into groups."""
    for s in shapes:
        if s.shape_id == shape_id:
            return s
        if s.shape_type == 6:  # GROUP
            r = find_shape_by_id(s.shapes, shape_id)
            if r is not None:
                return r
    return None


def set_textbox_text(shape, text: str, *, lines: list[str] | None = None) -> None:
    """Replace text inside a textbox shape, preserving the first run's formatting.

    If `lines` is given, render each as its own paragraph.
    If `text` is given, treat embedded newlines as paragraph breaks.
    """
    tf = shape.text_frame
    if lines is None:
        lines = text.split("\n")

    # Preserve formatting of the very first run (if any) as the template style.
    template_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            template_run = r
            break
        if template_run is not None:
            break

    # Capture template run XML to clone for new runs.
    template_rPr_xml = None
    if template_run is not None:
        rPr = template_run._r.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
        )
        if rPr is not None:
            template_rPr_xml = copy.deepcopy(rPr)

    # Capture first paragraph properties for cloning into new paragraphs.
    first_p = tf.paragraphs[0]
    pPr_xml = None
    pPr = first_p._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
    if pPr is not None:
        pPr_xml = copy.deepcopy(pPr)

    # Strategy: rewrite. Clear all paragraphs, then create fresh paragraphs
    # carrying the captured pPr/rPr.
    txBody = tf._txBody
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    # Remove all <a:p> children
    for p_el in list(txBody.findall(f"{a_ns}p")):
        txBody.remove(p_el)

    from lxml import etree

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for line in lines:
        p_el = etree.SubElement(txBody, f"{{{A}}}p")
        if pPr_xml is not None:
            p_el.append(copy.deepcopy(pPr_xml))
        r_el = etree.SubElement(p_el, f"{{{A}}}r")
        if template_rPr_xml is not None:
            r_el.append(copy.deepcopy(template_rPr_xml))
        else:
            etree.SubElement(r_el, f"{{{A}}}rPr", {"lang": "pt-PT"})
        t_el = etree.SubElement(r_el, f"{{{A}}}t")
        t_el.text = line if line else ""


def shrink_font(shape, size_pt: float) -> None:
    """Force-shrink all run font sizes inside a textbox."""
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for r in shape.text_frame._txBody.iter(f"{{{A}}}rPr"):
        r.set("sz", str(int(size_pt * 100)))


# ---------------------------------------------------------------------------
# Render orchestration
# ---------------------------------------------------------------------------


def build_evidence_table_lines(rows: list[dict]) -> list[str]:
    """One paragraph per row: 'EV-ID — aspecto — número (fonte)'."""
    out = []
    for r in rows:
        line = (
            f"{r['evidence_id']} — {r['aspect_pt']} — "
            f"{r['number']} ({r['source']})"
        )
        out.append(line)
    return out


def build_pilot_sketch_lines(p: dict) -> list[str]:
    """Compact pilot sketch summary for left box."""
    lines = []
    lines.append("Hipótese operacional:")
    lines.append(p["operational_hypothesis"])
    lines.append("")
    lines.append("Coorte: 40 estudantes 18–25, 1.º contacto digital SASU/linha municipal.")
    lines.append("Stack: mindLAMP self-hosted Hetzner DE (GDPR); EMA + JITAI rule-based; sem ML; sem wearables.")
    lines.append("Mediação: 1 psicólogo coordenador + 1 referenciador (~5h/sem). Encaminhamentos manuais.")
    lines.append("")
    lines.append("Cronograma (10 sem):")
    for sem, act, deliv in p["timeline_rows"]:
        lines.append(f"• Sem {sem}: {act}")
    lines.append("")
    lines.append("Endpoints: retenção 1.º→2.º contacto @4 sem (≥70%); compliance EMA ≥45%; Δ GAD-7/PHQ-9 ≥20% wk 8.")
    return lines


def render_pptx(parts: dict) -> tuple[Path, list[str]]:
    warnings: list[str] = []
    prs = Presentation(str(PPTX_TEMPLATE))
    slide1 = prs.slides[0]
    slide2 = prs.slides[1]

    def inner_textbox(group):
        """Pick the TextBox shape inside a group (the one carrying the
        '(preencher)' / '(principal)' placeholder run), not the rounded
        rectangle behind it."""
        # Prefer shapes whose existing text contains a placeholder marker.
        candidates = [
            s for s in group.shapes
            if s.has_text_frame and any(
                m in s.text_frame.text for m in ("(preencher)", "(principal)", "(suporte)")
            )
        ]
        if candidates:
            return candidates[0]
        # Fallback: prefer TEXT_BOX type (17) over AUTO_SHAPE.
        for s in group.shapes:
            if s.has_text_frame and s.shape_type == 17:
                return s
        return next(s for s in group.shapes if s.has_text_frame)

    # ---- Slide 1 ----
    # Problem Statement: Group 21 -> the inner TextBox shape
    g_problem = find_shape_by_id(slide1.shapes, 21)
    tb_problem = inner_textbox(g_problem)
    if len(parts["problem_text"]) > PROBLEM_MAX:
        warnings.append(
            f"Problem Statement length {len(parts['problem_text'])} > {PROBLEM_MAX}"
        )
    set_textbox_text(tb_problem, parts["problem_text"])

    # Hipótese de valor: Group 6
    g_hyp = find_shape_by_id(slide1.shapes, 6)
    tb_hyp = inner_textbox(g_hyp)
    hyp_text = parts["value"]["main_hypothesis"]
    if len(hyp_text) > HYPOTHESIS_MAX + 20:
        warnings.append(f"Main hypothesis length {len(hyp_text)} > soft limit")
    set_textbox_text(tb_hyp, hyp_text)

    # Métrica principal: Group 16 (single-line band)
    g_pm = find_shape_by_id(slide1.shapes, 16)
    tb_pm = inner_textbox(g_pm)
    pm_line = (
        f"{parts['value']['principal_metric_name']} — "
        f"baseline {parts['value']['principal_metric_baseline']}; "
        f"target {parts['value']['principal_metric_target']}"
    )
    if len(pm_line) > METRIC_LINE_MAX:
        warnings.append(f"Principal metric line {len(pm_line)} > {METRIC_LINE_MAX}")
    set_textbox_text(tb_pm, pm_line)

    # Suporte 1 / 2: TextBox 35 / 36
    tb_s1 = find_shape_by_id(slide1.shapes, 35)
    tb_s2 = find_shape_by_id(slide1.shapes, 36)
    s1_line = f"{parts['value']['support1_name']} — {parts['value']['support1_hypothesis']}"
    s2_line = f"{parts['value']['support2_name']} — {parts['value']['support2_hypothesis']}"
    for label, line in (("support1", s1_line), ("support2", s2_line)):
        if len(line) > SUPPORT_LINE_MAX:
            warnings.append(f"{label} line {len(line)} > {SUPPORT_LINE_MAX}")
    set_textbox_text(tb_s1, s1_line)
    set_textbox_text(tb_s2, s2_line)

    # Tabela de Evidências: Group 4 (Slide 1) -> textbox lines
    g_evid = find_shape_by_id(slide1.shapes, 4)
    tb_evid = inner_textbox(g_evid)
    evid_lines = build_evidence_table_lines(parts["evidence_table"]["rows"])
    for line in evid_lines:
        if len(line) > EVIDENCE_ROW_MAX:
            warnings.append(f"Evidence row {len(line)} > {EVIDENCE_ROW_MAX}: {line[:60]}…")
    set_textbox_text(tb_evid, "\n".join(evid_lines), lines=evid_lines)
    shrink_font(tb_evid, 10.0)

    # ---- Slide 2 ----
    # Pilot Sketch: Group 4 (in slide2)
    g_pilot = find_shape_by_id(slide2.shapes, 4)
    tb_pilot = inner_textbox(g_pilot)
    pilot_lines = build_pilot_sketch_lines(parts["pilot"])
    pilot_text_total = "\n".join(pilot_lines)
    if len(pilot_text_total) > PILOT_SKETCH_MAX:
        warnings.append(
            f"Pilot sketch total length {len(pilot_text_total)} > {PILOT_SKETCH_MAX}"
        )
    set_textbox_text(tb_pilot, pilot_text_total, lines=pilot_lines)
    shrink_font(tb_pilot, 10.0)

    # Budget Range: Group 20 (single-line)
    g_budget = find_shape_by_id(slide2.shapes, 20)
    tb_budget = inner_textbox(g_budget)
    budget_line = parts["pilot"]["budget_total"]
    if len(budget_line) > BUDGET_LINE_MAX:
        warnings.append(f"Budget line {len(budget_line)} > {BUDGET_LINE_MAX}")
    set_textbox_text(tb_budget, budget_line)

    # Budget notes (Group 25 right below Budget Range)
    g_budget_notes = find_shape_by_id(slide2.shapes, 25)
    tb_budget_notes = inner_textbox(g_budget_notes)
    bn = parts["pilot"]["budget_notes"]
    if len(bn) > BUDGET_NOTES_MAX:
        warnings.append(f"Budget notes {len(bn)} > {BUDGET_NOTES_MAX}")
    set_textbox_text(tb_budget_notes, bn)
    shrink_font(tb_budget_notes, 10.0)

    # Pressupostos: TextBox 39, 40 (top 2 critical assumptions)
    tb_a1 = find_shape_by_id(slide2.shapes, 39)
    tb_a2 = find_shape_by_id(slide2.shapes, 40)
    assumptions = parts["pilot"]["assumptions"]
    if len(assumptions) < 2:
        raise ValueError("Need at least 2 critical assumptions")

    # Strip the leading "**Title**:" markdown bold and keep concise text.
    def clean_assumption(a: str) -> str:
        a = re.sub(r"\*\*", "", a)
        return a

    a1 = clean_assumption(assumptions[0])
    a2 = clean_assumption(assumptions[1])
    for label, line in (("assumption1", a1), ("assumption2", a2)):
        if len(line) > ASSUMPTION_MAX:
            warnings.append(f"{label} length {len(line)} > {ASSUMPTION_MAX}")
    set_textbox_text(tb_a1, a1)
    set_textbox_text(tb_a2, a2)
    shrink_font(tb_a1, 9.0)
    shrink_font(tb_a2, 9.0)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT, warnings


# ---------------------------------------------------------------------------
# DOCX rendering
# ---------------------------------------------------------------------------


def render_docx(parts: dict, evids_in_deck: list[str]) -> tuple[Path, list[str]]:
    warnings: list[str] = []
    doc = Document(str(DOCX_TEMPLATE))

    # Replace the "(preencher)" lines in P[4]
    target_p = None
    for p in doc.paragraphs:
        if "(preencher)" in p.text:
            target_p = p
            break
    if target_p is None:
        raise RuntimeError("Could not find (preencher) paragraph in docx template")

    new_text = (
        f"Nome do projeto/startup: {PROJECT_NAME}\n"
        f"Pessoa de contacto (nome completo e contacto): {CONTACT}"
    )

    # Replace each run; clear all but first then put text into first.
    for i, run in enumerate(target_p.runs):
        run.text = "" if i > 0 else new_text
    # If somehow no runs, add one.
    if not target_p.runs:
        target_p.add_run(new_text)

    # Append "Notas Adicionais" body content.
    ledger_by_id = {e["id"]: e for e in parts["ledger"]["entries"]}

    def add_h(text: str, size: int = 14, bold: bool = True) -> None:
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = bold
        r.font.size = DocxPt(size)

    def add_p(text: str, size: int = 10, bold: bool = False) -> None:
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = bold
        r.font.size = DocxPt(size)

    add_h("Sources scanned", size=12)
    for s in parts["ledger"]["metadata"]["sources_scanned"]:
        add_p(f"• {s}")
    skipped = parts["ledger"]["metadata"].get("sources_skipped", [])
    if skipped:
        add_p("Sources skipped:", bold=True)
        for s in skipped:
            add_p(f"• {s['file']} — {s['reason']}")

    add_h("Methodology", size=12)
    add_p(
        "Quantitative claims foram extraídos automaticamente dos PDFs e "
        "HTMLs do dossier (scripts/extract_text.py + find_claims.py + "
        "build_ledger.py), atribuídos um EV-ID estável (EV-001…EV-088), "
        "categorizados por tópico (prevalence, access, user_segment, "
        "digital_health, platform, pilot_metric, cost, other) e normalizados "
        "para um JSON ledger. Cada claim citado nas peças (problem, value, "
        "evidence_table, pilot) referencia o seu EV-ID; este documento lista "
        "as notas para os IDs efectivamente usados na deck."
    )

    add_h("Per-claim notes (EV-IDs cited in deck)", size=12)
    documented_ids: list[str] = []
    for evid in evids_in_deck:
        entry = ledger_by_id.get(evid)
        add_p(evid, bold=True, size=11)
        if entry is None:
            add_p("⚠ Not present in ledger — investigate.", size=10)
            continue
        documented_ids.append(evid)
        add_p(f"Claim: {entry['claim_pt']}")
        meta = (
            f"Number: {entry.get('number') or '—'}  |  "
            f"Metric: {entry.get('metric') or '—'}  |  "
            f"Topic: {entry.get('topic') or '—'}  |  "
            f"Confidence: {entry.get('confidence') or '—'}"
        )
        add_p(meta, size=9)
        src = (
            f"Source: {entry.get('source_file') or '—'}"
            + (f" — {entry['source_section']}" if entry.get("source_section") else "")
            + (f" (p.{entry['source_page']})" if entry.get("source_page") else "")
        )
        add_p(src, size=9)
        if entry.get("original_context"):
            add_p(f"Original context: \"{entry['original_context']}\"", size=9)
        if entry.get("notes"):
            add_p(f"Notes: {entry['notes']}", size=9)

    add_h("Coverage check", size=12)
    add_p(
        f"EV-IDs cited in deck: {len(evids_in_deck)} — {', '.join(evids_in_deck)}"
    )
    add_p(
        f"EV-IDs documented in this notes file: {len(documented_ids)} — "
        f"{', '.join(documented_ids)}"
    )
    missing = sorted(set(evids_in_deck) - set(documented_ids))
    if missing:
        add_p(f"⚠ Missing from notes: {', '.join(missing)}", bold=True)
        warnings.append(f"DOCX missing EV-IDs: {missing}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(DOCX_OUT))
    return DOCX_OUT, warnings


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------


def count_placeholders(path: Path, patterns: list[str]) -> dict[str, int]:
    counts = {p: 0 for p in patterns}
    with zipfile.ZipFile(path, "r") as z:
        for name in z.namelist():
            if not (name.endswith(".xml") or name.endswith(".rels")):
                continue
            data = z.read(name).decode("utf-8", errors="ignore")
            for pat in patterns:
                counts[pat] += data.count(pat)
    return counts


def evids_in_pptx(path: Path) -> set[str]:
    found: set[str] = set()
    with zipfile.ZipFile(path, "r") as z:
        for name in z.namelist():
            if "/slides/slide" in name and name.endswith(".xml"):
                data = z.read(name).decode("utf-8", errors="ignore")
                # Strip XML tags so EV-IDs split across runs still match.
                stripped = re.sub(r"<[^>]+>", "", data)
                found.update(EV_ID_RE.findall(stripped))
    return found


def evids_in_docx(path: Path) -> set[str]:
    found: set[str] = set()
    with zipfile.ZipFile(path, "r") as z:
        for name in z.namelist():
            if name.endswith(".xml"):
                data = z.read(name).decode("utf-8", errors="ignore")
                stripped = re.sub(r"<[^>]+>", "", data)
                found.update(EV_ID_RE.findall(stripped))
    return found


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    drafts = load_drafts()
    parts = {
        "problem_text": extract_problem(drafts["problem_md"]),
        "value": extract_value(drafts["value_md"]),
        "pilot": extract_pilot(drafts["pilot_md"]),
        "evidence_table": drafts["evidence_table"],
        "ledger": drafts["ledger"],
    }

    deck_evids = collect_evids_in_deck_text(parts, parts["evidence_table"]["rows"])

    pptx_path, pptx_warnings = render_pptx(parts)
    docx_path, docx_warnings = render_docx(parts, deck_evids)

    # Validation
    placeholder_patterns = ["(preencher)", "(principal)", "(suporte)", "[…]", "[...]"]
    pptx_ph = count_placeholders(pptx_path, placeholder_patterns)
    docx_ph = count_placeholders(docx_path, placeholder_patterns)

    pptx_size = pptx_path.stat().st_size
    docx_size = docx_path.stat().st_size

    pptx_evids = evids_in_pptx(pptx_path)
    docx_evids = evids_in_docx(docx_path)

    coverage_missing = sorted(pptx_evids - docx_evids)

    # Hard checks
    failures = []
    if pptx_size < 20_000:
        failures.append(f"pptx size {pptx_size} < 20KB")
    if docx_size < 10_000:
        failures.append(f"docx size {docx_size} < 10KB")
    if sum(pptx_ph.values()) > 0:
        failures.append(f"pptx still has placeholders: {pptx_ph}")
    if sum(docx_ph.values()) > 0:
        failures.append(f"docx still has placeholders: {docx_ph}")
    if coverage_missing:
        failures.append(f"EV-IDs in deck missing from notes: {coverage_missing}")

    print("=" * 60)
    print("RENDER SUMMARY")
    print(
        f"  pptx: {pptx_path} ({pptx_size//1024} KB, "
        f"placeholders remaining: {sum(pptx_ph.values())})"
    )
    print(
        f"  docx: {docx_path} ({docx_size//1024} KB, "
        f"placeholders remaining: {sum(docx_ph.values())})"
    )
    print(f"  EV-IDs cited in deck: {sorted(pptx_evids, key=lambda x: int(x.split('-')[1]))}")
    print(
        f"  EV-IDs documented in notes: "
        f"{sorted(docx_evids, key=lambda x: int(x.split('-')[1]))}"
    )
    coverage_pct = (
        100.0 if not pptx_evids
        else 100.0 * len(pptx_evids & docx_evids) / len(pptx_evids)
    )
    print(f"  Coverage: {coverage_pct:.0f}%")
    if pptx_warnings or docx_warnings:
        print("  Warnings:")
        for w in pptx_warnings + docx_warnings:
            print(f"    - {w}")
    if failures:
        print("  FAILURES:")
        for f in failures:
            print(f"    - {f}")
        print("=" * 60)
        return 1
    print("=" * 60)
    return 0


if __name__ == "__main__":
    sys.exit(main())
